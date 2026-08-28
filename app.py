"""
Microservizio per la compilazione del template pptx (Valutazione Preliminare).
Gestisce:
  1. Sostituzione placeholder di testo (WWW.DOMINIOCLIENTE.IT, {{TOKEN}})
  2. Colore dinamico dei 4 cerchi punteggio PageSpeed (slide 3) e dei bordi
     Load time (slide 4), individuati tramite il campo "descr" (testo
     alternativo impostato in Google Slides)
  3. Analisi on-page gratuita (title/description mancanti, troppo corti/
     lunghi, duplicati) come sostituto minimale di Screaming Frog

Endpoint:
  POST /compila         -> compila il template (invariato)
  POST /analizza-onpage -> crawla un sito e restituisce errori title/description
  GET  /health
"""

import concurrent.futures
import io
import json
import re
import shutil
import tempfile
from urllib.parse import urljoin, urlparse

import requests
from bs4 import BeautifulSoup
from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.dml.color import RGBColor

app = FastAPI(title="MP Quadro - Compilazione template Valutazione Preliminare")

NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

COLORE_ROSSO = RGBColor(0xEA, 0x43, 0x35)
COLORE_GIALLO = RGBColor(0xF4, 0xB4, 0x00)
COLORE_VERDE = RGBColor(0x34, 0xA8, 0x53)
COLORE_GRIGIO = RGBColor(0x9E, 0x9E, 0x9E)  # dato mancante/errore


def colore_per_punteggio(punteggio):
    if punteggio is None:
        return COLORE_GRIGIO
    if punteggio >= 90:
        return COLORE_VERDE
    if punteggio >= 50:
        return COLORE_GIALLO
    return COLORE_ROSSO


def colore_per_secondi(secondi):
    if secondi is None:
        return COLORE_GRIGIO
    if secondi <= 1.5:
        return COLORE_VERDE
    if secondi <= 3:
        return COLORE_GIALLO
    return COLORE_ROSSO


def colora_testo_livello(text_frame):
    """Colora il run il cui testo finale è esattamente OTTIMO/BUONO/PESSIMO/N/D."""
    mappa = {"OTTIMO": COLORE_VERDE, "BUONO": COLORE_GIALLO, "PESSIMO": COLORE_ROSSO, "N/D": COLORE_GRIGIO}
    n = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() in mappa:
                run.font.color.rgb = mappa[run.text.strip()]
                n += 1
    return n


def livello_per_tempo(secondi):
    if secondi is None:
        return "N/D"
    if secondi <= 1.5:
        return "OTTIMO"
    if secondi <= 3:
        return "BUONO"
    return "PESSIMO"


def costruisci_sostituzioni(dati):
    """Costruisce il dizionario {placeholder: valore} a partire dal payload n8n."""
    pagespeed = {(r["label"], r["strategy"]): r for r in dati.get("pagespeed", [])}

    def val(label, strategy, campo, default="N/D"):
        r = pagespeed.get((label, strategy))
        if r is None or r.get("errore"):
            return default
        v = r.get(campo)
        return default if v is None else str(v)

    cliente_desktop = pagespeed.get(("cliente", "desktop"))
    competitor_desktop = pagespeed.get(("competitor", "desktop"))

    onpage = {r["label"]: r for r in dati.get("onpage", [])}

    def val_onpage(label, campo, default="N/D"):
        r = onpage.get(label)
        if r is None or r.get("errore"):
            return default
        v = r.get(campo)
        return default if v is None else str(v)

    sostituzioni = {
        "PSI_MOBILE_CLIENTE": val("cliente", "mobile", "punteggio"),
        "PSI_DESKTOP_CLIENTE": val("cliente", "desktop", "punteggio"),
        "PSI_MOBILE_COMPETITOR": val("competitor", "mobile", "punteggio"),
        "PSI_DESKTOP_COMPETITOR": val("competitor", "desktop", "punteggio"),
        "GRADO_CLIENTE": val("cliente", "desktop", "grado"),
        "GRADO_COMPETITOR": val("competitor", "desktop", "grado"),
        "PAGESIZE_CLIENTE": val("cliente", "desktop", "pageSizeMB"),
        "PAGESIZE_COMPETITOR": val("competitor", "desktop", "pageSizeMB"),
        "LOADTIME_CLIENTE": val("cliente", "desktop", "loadTimeSec"),
        "LOADTIME_COMPETITOR": val("competitor", "desktop", "loadTimeSec"),
        "REQUESTS_CLIENTE": val("cliente", "desktop", "numeroRichieste"),
        "REQUESTS_COMPETITOR": val("competitor", "desktop", "numeroRichieste"),
        "LIVELLO_CLIENTE": livello_per_tempo(
            cliente_desktop.get("loadTimeSec") if cliente_desktop and not cliente_desktop.get("errore") else None
        ),
        "LIVELLO_COMPETITOR": livello_per_tempo(
            competitor_desktop.get("loadTimeSec") if competitor_desktop and not competitor_desktop.get("errore") else None
        ),
        "ERRORI_TITLE_CLIENTE": val_onpage("cliente", "erroriTitle"),
        "ERRORI_TITLE_COMPETITOR": val_onpage("competitor", "erroriTitle"),
        "ERRORI_DESCRIPTION_CLIENTE": val_onpage("cliente", "erroriDescription"),
        "ERRORI_DESCRIPTION_COMPETITOR": val_onpage("competitor", "erroriDescription"),
    }

    return sostituzioni, pagespeed


def sostituisci_in_run(run, mappa):
    cambiato = False
    testo = run.text

    for placeholder, valore in [
        ("WWW.DOMINIOCLIENTE.IT", mappa.get("_nomeCliente")),
        ("WWW.DOMINIOCOMPETITOR.IT", mappa.get("_nomeCompetitor")),
    ]:
        if valore and placeholder.lower() in testo.lower():
            pattern = re.compile(re.escape(placeholder), re.IGNORECASE)
            testo = pattern.sub(valore, testo)
            cambiato = True

    for token, valore in mappa.items():
        if token.startswith("_"):
            continue
        marcatore = "{{" + token + "}}"
        if marcatore in testo:
            testo = testo.replace(marcatore, valore)
            cambiato = True

    if cambiato:
        run.text = testo
    return cambiato


def processa_text_frame(text_frame, mappa):
    n = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if sostituisci_in_run(run, mappa):
                n += 1
    return n


def processa_tabella(table, mappa):
    n = 0
    for row in table.rows:
        for cell in row.cells:
            n += processa_text_frame(cell.text_frame, mappa)
    return n


def trova_descr(shape):
    el = shape._element.find(f".//{NS_P}cNvPr")
    return el.get("descr") if el is not None else None


def colora_cerchi(slide, pagespeed):
    mappa_forme = {
        "circle_psi_mobile_cliente": ("cliente", "mobile"),
        "circle_psi_desktop_cliente": ("cliente", "desktop"),
        "circle_psi_mobile_competitor": ("competitor", "mobile"),
        "circle_psi_desktop_competitor": ("competitor", "desktop"),
    }
    mappa_box_loadtime = {
        "loadtime_box_cliente": "cliente",
        "loadtime_box_competitor": "competitor",
    }
    colorati = 0

    def ricorsione(shapes):
        nonlocal colorati
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                ricorsione(shape.shapes)
                continue
            d = trova_descr(shape)
            if d in mappa_forme:
                label, strategy = mappa_forme[d]
                r = pagespeed.get((label, strategy))
                punteggio = r.get("punteggio") if r and not r.get("errore") else None
                colore = colore_per_punteggio(punteggio)
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = colore
                    colorati += 1
                except Exception:
                    pass
            elif d in mappa_box_loadtime:
                label = mappa_box_loadtime[d]
                r = pagespeed.get((label, "desktop"))
                secondi = r.get("loadTimeSec") if r and not r.get("errore") else None
                colore = colore_per_secondi(secondi)
                try:
                    shape.line.color.rgb = colore
                    colorati += 1
                except Exception:
                    pass

    ricorsione(slide.shapes)
    return colorati


@app.post("/compila")
def compila(
    file: UploadFile = File(...),
    dati: str = Form(...),
):
    mime_pptx = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    nome_ok = (file.filename or "").endswith(".pptx")
    mime_ok = file.content_type == mime_pptx
    if not nome_ok and not mime_ok:
        raise HTTPException(
            400,
            f"Il file deve essere .pptx (ricevuto filename={file.filename!r}, "
            f"content_type={file.content_type!r})",
        )

    try:
        payload = json.loads(dati)
    except json.JSONDecodeError:
        raise HTTPException(400, "Campo 'dati' non è JSON valido")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        shutil.copyfileobj(file.file, tmp)
        tmp_path = tmp.name

    prs = Presentation(tmp_path)

    sostituzioni, pagespeed = costruisci_sostituzioni(payload)
    sostituzioni["_nomeCliente"] = payload.get("nomeCliente")
    sostituzioni["_nomeCompetitor"] = payload.get("nomeCompetitor")

    totale_testo = 0
    totale_colori = 0
    totale_livello = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                totale_testo += processa_text_frame(shape.text_frame, sostituzioni)
                totale_livello += colora_testo_livello(shape.text_frame)
            if shape.has_table:
                totale_testo += processa_tabella(shape.table, sostituzioni)
        totale_colori += colora_cerchi(slide, pagespeed)

    if totale_testo == 0:
        raise HTTPException(
            422,
            "Nessuna sostituzione di testo effettuata. Verifica che il template "
            "contenga ancora i placeholder attesi.",
        )

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": 'attachment; filename="output.pptx"',
            "X-Sostituzioni-Testo": str(totale_testo),
            "X-Cerchi-Colorati": str(totale_colori),
        },
    )


@app.get("/health")
async def health():
    return {"status": "ok"}


# ---------------------------------------------------------------------------
# Analisi on-page gratuita (sostituto minimale di Screaming Frog)
# ---------------------------------------------------------------------------

USER_AGENT = "Mozilla/5.0 (compatible; MPQuadroValutazioneBot/1.0)"
TIMEOUT_SEC = 10


ESTENSIONI_NON_PAGINA = {
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg", ".bmp", ".ico",
    ".pdf", ".zip", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".mp4", ".mp3", ".avi", ".mov", ".css", ".js", ".xml", ".json",
}

# Pattern di percorso tipici di pagine generate automaticamente da CMS
# (archivi categoria/tag, paginazione, feed, autore) che normalmente non
# hanno title/description propri e che un audit SEO reale esclude di default
# (Screaming Frog e simili li escludono o li tratta come categoria separata).
PATTERN_NON_CONTENUTO = [
    "/category/", "/categoria/", "/tag/", "/tags/", "/author/", "/autore/",
    "/page/", "/pagina/", "/feed/", "/wp-json/", "/wp-admin/", "/wp-login",
    "/attachment/",
]


def normalizza_per_confronto(url):
    """Riduce un URL alla sua forma essenziale (dominio + percorso) ignorando
    http/https, www, e slash finale — per riconoscere come 'stessa pagina'
    varianti come https://www.sito.it e http://sito.it/ che altrimenti
    verrebbero contate come pagine distinte e segnalate come falsi duplicati."""
    p = urlparse(url)
    netloc = p.netloc.lower()
    if netloc.startswith("www."):
        netloc = netloc[4:]
    path = p.path.rstrip("/") or "/"
    return netloc + path


def e_una_pagina(url):
    """Esclude file media/documenti/asset tecnici e pagine CMS generate automaticamente
    (categorie, tag, paginazione) che non sono contenuto editoriale vero e proprio."""
    path = urlparse(url).path.lower()
    if any(path.endswith(ext) for ext in ESTENSIONI_NON_PAGINA):
        return False
    if any(pattern in path for pattern in PATTERN_NON_CONTENUTO):
        return False
    return True


def stesso_dominio(url, dominio_base):
    try:
        return urlparse(url).netloc.lower().lstrip("www.") == dominio_base.lower().lstrip("www.")
    except Exception:
        return False


def scopri_pagine(base_url, max_pagine=25):
    """Prova prima la sitemap.xml (gestendo anche sitemap-indice); se assente/vuota, fa un crawl BFS semplice."""
    dominio = urlparse(base_url).netloc
    pagine = []

    try:
        r = requests.get(urljoin(base_url, "/sitemap.xml"), headers={"User-Agent": USER_AGENT}, timeout=TIMEOUT_SEC)
        if r.status_code == 200 and "xml" in r.headers.get("Content-Type", ""):
            soup = BeautifulSoup(r.content, "xml")

            if soup.find("sitemapindex"):
                # È un indice: contiene link ad altre sitemap, non pagine dirette
                sotto_sitemap = [loc.text.strip() for loc in soup.find_all("loc")][:10]
                for sm_url in sotto_sitemap:
                    try:
                        r2 = requests.get(sm_url, headers={"User-Agent": USER_AGENT}, timeout=TIMEOUT_SEC)
                        if r2.status_code != 200:
                            continue
                        soup2 = BeautifulSoup(r2.content, "xml")
                        locs2 = [loc.text.strip() for loc in soup2.find_all("loc")]
                        pagine.extend([u for u in locs2 if stesso_dominio(u, dominio) and e_una_pagina(u)])
                    except Exception:
                        continue
            else:
                # Sitemap diretta: <urlset> con pagine vere
                locs = [loc.text.strip() for loc in soup.find_all("loc")]
                pagine = [u for u in locs if stesso_dominio(u, dominio) and e_una_pagina(u)]
    except Exception:
        pass

    # Deduplica PRIMA di tagliare a max_pagine: un URL può comparire in più
    # sotto-sitemap (capita su WordPress), e senza questo passaggio verrebbe
    # scaricato e valutato più volte, facendo scattare falsi "duplicato"
    # contro se stesso invece che contro un'altra pagina reale.
    # Deduplica PRIMA di tagliare a max_pagine: la stessa pagina può comparire
    # più volte con varianti diverse (con/senza www, http/https, slash finale,
    # o ripetuta in più sotto-sitemap) — senza normalizzare per il confronto,
    # verrebbe scaricata più volte e segnalata come falso "duplicato" contro
    # se stessa invece che contro un'altra pagina reale.
    visti = set()
    pagine_dedup = []
    for u in pagine:
        chiave = normalizza_per_confronto(u)
        if chiave not in visti:
            visti.add(chiave)
            pagine_dedup.append(u)
    pagine = pagine_dedup[:max_pagine]

    if pagine:
        return pagine

    # Tentativo 2: crawl BFS semplice partendo dalla homepage
    visitate = set()
    da_visitare = [base_url]
    while da_visitare and len(visitate) < max_pagine:
        url = da_visitare.pop(0)
        if url in visitate:
            continue
        visitate.add(url)
        try:
            r = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=TIMEOUT_SEC)
            if r.status_code != 200:
                continue
            soup = BeautifulSoup(r.content, "html.parser")
            for a in soup.find_all("a", href=True):
                link = urljoin(url, a["href"]).split("#")[0]
                if stesso_dominio(link, dominio) and e_una_pagina(link) and link not in visitate and link not in da_visitare:
                    da_visitare.append(link)
        except Exception:
            continue

    return list(visitate)[:max_pagine]


def analizza_pagina(url):
    try:
        r = requests.get(url, headers={"User-Agent": USER_AGENT}, timeout=TIMEOUT_SEC, allow_redirects=True)
    except Exception as e:
        return {"url": url, "status": None, "errore_richiesta": str(e), "title": None, "description": None}

    risultato = {"url": url, "status": r.status_code, "errore_richiesta": None}

    if r.status_code >= 400:
        risultato["title"] = None
        risultato["description"] = None
        return risultato

    soup = BeautifulSoup(r.content, "html.parser")
    title_tag = soup.find("title")
    desc_tag = soup.find("meta", attrs={"name": "description"})

    risultato["title"] = title_tag.text.strip() if title_tag and title_tag.text else None
    risultato["description"] = desc_tag.get("content", "").strip() if desc_tag and desc_tag.get("content") else None

    return risultato


def valuta_title(title):
    if not title:
        return "mancante"
    return None


def valuta_description(desc):
    if not desc:
        return "mancante"
    return None


def valuta_risultati_pagine(risultati, dominio):
    """Aggrega i risultati di crawl di un singolo sito in un riepilogo errori."""
    titoli = [r["title"] for r in risultati if r["title"]]
    descrizioni = [r["description"] for r in risultati if r["description"]]
    duplicati_titoli = {t for t in titoli if titoli.count(t) > 1}
    duplicati_descrizioni = {d for d in descrizioni if descrizioni.count(d) > 1}

    pagine_con_errore_title = []
    pagine_con_errore_description = []
    pagine_con_errore_status = []

    for r in risultati:
        if r["status"] is None or r["status"] >= 400:
            pagine_con_errore_status.append({"url": r["url"], "status": r["status"], "dettaglio": r.get("errore_richiesta")})
            continue
        problema_title = valuta_title(r["title"])
        if not problema_title and r["title"] in duplicati_titoli:
            problema_title = "duplicato"
        if problema_title:
            pagine_con_errore_title.append({"url": r["url"], "problema": problema_title, "valore": r["title"]})

        problema_desc = valuta_description(r["description"])
        if not problema_desc and r["description"] in duplicati_descrizioni:
            problema_desc = "duplicata"
        if problema_desc:
            pagine_con_errore_description.append({"url": r["url"], "problema": problema_desc, "valore": r["description"]})

    return {
        "dominio": dominio,
        "pagine_analizzate": len(risultati),
        "conteggio_errori_title": len(pagine_con_errore_title),
        "conteggio_errori_description": len(pagine_con_errore_description),
        "conteggio_errori_status": len(pagine_con_errore_status),
        "dettaglio_errori_title": pagine_con_errore_title,
        "dettaglio_errori_description": pagine_con_errore_description,
        "dettaglio_errori_status": pagine_con_errore_status,
    }


@app.post("/analizza-onpage")
def analizza_onpage(url: str = Form(...), max_pagine: int = Form(25)):
    if not url.startswith(("http://", "https://")):
        url = "https://" + url

    max_pagine = min(max_pagine, 200)

    pagine = scopri_pagine(url, max_pagine)
    if not pagine:
        raise HTTPException(422, "Nessuna pagina raggiungibile per questo dominio")

    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        risultati = list(executor.map(analizza_pagina, pagine))

    risultato = valuta_risultati_pagine(risultati, url)
    risultato["nota"] = (
        "Analisi indicativa: copertura limitata a max_pagine, non esegue JavaScript "
        "(possibili falsi positivi su siti che generano title/description via JS). "
        "Segnala solo title/description mancanti o duplicati tra le pagine analizzate."
    )
    return risultato


@app.post("/analizza-onpage-duplice")
def analizza_onpage_duplice(
    url_cliente: str = Form(...),
    url_competitor: str = Form(...),
    max_pagine: int = Form(25),
):
    """
    Analizza cliente e competitor in UNA sola chiamata, condividendo lo stesso
    pool di thread: evita che n8n, chiamando i due siti in due richieste HTTP
    separate (spesso in sequenza, non in parallelo), sommi i tempi dei due
    crawl invece di sovrapporli.
    """

    def prepara(url):
        if not url.startswith(("http://", "https://")):
            url = "https://" + url
        return url

    url_cliente = prepara(url_cliente)
    url_competitor = prepara(url_competitor)
    max_pagine = min(max_pagine, 200)

    pagine_cliente = scopri_pagine(url_cliente, max_pagine)
    pagine_competitor = scopri_pagine(url_competitor, max_pagine)

    if not pagine_cliente and not pagine_competitor:
        raise HTTPException(422, "Nessuna pagina raggiungibile né per il cliente né per il competitor")

    # Stesso pool condiviso per entrambi i siti: le richieste dei due domini
    # si sovrappongono nel tempo invece di essere l'una in coda all'altra.
    with concurrent.futures.ThreadPoolExecutor(max_workers=15) as executor:
        future_cliente = [executor.submit(analizza_pagina, p) for p in pagine_cliente]
        future_competitor = [executor.submit(analizza_pagina, p) for p in pagine_competitor]
        risultati_cliente = [f.result() for f in future_cliente]
        risultati_competitor = [f.result() for f in future_competitor]

    return {
        "cliente": valuta_risultati_pagine(risultati_cliente, url_cliente) if pagine_cliente else {
            "dominio": url_cliente, "pagine_analizzate": 0, "conteggio_errori_title": None,
            "conteggio_errori_description": None, "conteggio_errori_status": None,
        },
        "competitor": valuta_risultati_pagine(risultati_competitor, url_competitor) if pagine_competitor else {
            "dominio": url_competitor, "pagine_analizzate": 0, "conteggio_errori_title": None,
            "conteggio_errori_description": None, "conteggio_errori_status": None,
        },
        "nota": (
            "Analisi indicativa: copertura limitata a max_pagine, non esegue JavaScript. "
            "Segnala solo title/description mancanti o duplicati tra le pagine analizzate "
            "(duplicati calcolati separatamente per ciascun dominio)."
        ),
    }
